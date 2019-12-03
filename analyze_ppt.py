from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches, Pt
import collections
import json
import argparse
import imageio
import os


def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Analyze powerpoint file structure')
    parser.add_argument('infile',
                        type=argparse.FileType('r'),
                        help='Powerpoint file to be analyzed')
    parser.add_argument('outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint')
    return parser.parse_args()


def analyze_ppt(input, output):
    """ Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    """
    # JSON String
    data = """
    {
        "Programming languages": {
            "": {
                "item1": "Can be thought of as simply a program that you feed commands to",
                "item2": "You can write scripts or programs that make the computer do whatever you can dream up",
                "item3": "There are many programming languages at least 700 or so according to Wikipedia",
                "image1": {
                    "slide1.jpg": "bottom"
                }
            }
        },
        "Brief history of programming languages": {
            "Aaron's list of most popular languages": {
                "item1": "1957 - FORTRAN - Compiled",
                "item2": "1964 - BASIC - Interpreted",
                "item3": "1970 - Pascal - Compiled",
                "item4": "1972 - C - Compiled",
                "item5": "1980 - C++ - Compiled",
                "item7": "1991 - Python - Interpreted",
                "item8": "1991 - Visual Basic - Compiled",
                "item10": "1995 - Ruby - Interpreted",
                "item11": "1995 - Java - Compiled (JVM)",
                "item12": "1995 - JavaScript - Interpreted",
                "item13": "1995 - PHP - Interpreted",
                "item14": "2001 - C# - Compiled (CLR)",
                "item15": "2009 - Go - Compiled (Google - produces statically linked native binaries without external dependencies.)",
                "item16": "2011 - Dart Compiled/Interpreted (Google - AOT-compiled to JavaScript)",
                "image1": {
                    "slide2.jpg": "right"
                }
            },
            "slide_layout": 2
        },
    	"Interpreted vs Compiled": {
            "": {
                "item1": "Compiled languages - converted directly into machine code that the processor can execute. As a result, they tend to be faster and more efficient to execute than interpreted languages.",
                "item2": "Interpreted languages - the source code is not directly translated by the target machine. Instead, a different program, aka the interpreter, reads and executes the code.",
                "image1": {
                    "slide3.jpg": "bottom"
                }
            },
            "slide_layout": 2
        },
        "Scripts vs Programs": {
            "": {
                "item1": "Scripts are usually interpreted (but not always, such as with golang scripts)",
                "item2": "Programs can be either compiled (C++, C#, Java) or interpreted (Python, Ruby, PHP)",
                "item3": "The biggest difference is that scripts are written to control an existing program",
                "item4": "Scripts often automate manual tasks to make work easier to accomplish",
                "item5": "Scripts can accomplish many important tasks and are often written by a single person",
                "item6": "Programs usually have very ambitious goals and often take a large amount of time and money to create",
                "image1": {
                    "slide4.jpg": "bottom_right"
                }
            },
            "slide_layout": 2
        },
        "Common Python Data types": {
            "Python is a dynamically typed language which means variables themselves are not bound to a specific data type. That said the following are the most commonly used data types in Python": {
                "item1": "Integers (123)",
                "item2": "Strings (abc)",
                "item3": "Boolean (True/False)",
                "item4": "Lists [1, 2, 3]",
                "item5": "Dictionaries {Key: Value}",
                "image": {
                    "slide5.jpg": "bottom_right"
                }
            },
            "slide_layout": 2
        },
        "Integers": {
            "Integers in python are positive or negative whole numbers with no decimal point": {
                "item1": "Python 2.7.16 (default, Oct  7 2019, 17:36:04) \n[GCC 8.3.0] on linux2\nType \\"help\\", \\"copyright\\", \\"credits\\" or \\"license\\" for more information.\n>>> x = 1\n>>> y = 2\n>>> type(x)\n<type 'int'>\n>>> type(y)\n<type 'int'>\n>>> print(x + y)\n3",
                "image": {
                    "slide6.jpg": "bottom_right"
                }
            },
            "slide_layout": 2
        },
        "Strings": {
            "In python a string is most easily identified by the use of double quotes": {
                "item1": "Python 2.7.16 (default, Oct  7 2019, 17:36:04) \n[GCC 8.3.0] on linux2\nType \\"help\\", \\"copyright\\", \\"credits\\" or \\"license\\" for more information.\n>>> x = \\"This is a string\\"\n>>> type(x)\n<type 'str'>\n>>> print(x)\nThis is a string",
                "image": {
                    "slide7.jpg": "bottom"
                }
            },
            "slide_layout": 2
        },
        "Boolean": {
            "Boolean simply means True or False": {
                "item1": "Python 2.7.16 (default, Oct  7 2019, 17:36:04) \n[GCC 8.3.0] on linux2\nType \\"help\\", \\"copyright\\", \\"credits\\" or \\"license\\" for more information.\n>>> x = True\n>>> type(x)\n<type 'bool'>\n>>> if x is True:\n...     print(\\"Boolean is pronounced boo-lee-uhn\\")\n... else:\n...     print(\\"Boolean is pronounced bool-yaan\\")\n... \nBoolean is pronounced boo-lee-uhn",
                "image": {
                    "slide8.jpg": "bottom_right"
                }
            },
            "slide_layout": 2
        },
        "Lists": {
            "Lists are denoted by square brackets [ ] and contain comma separated values. Another name for a list is an Array.\nIf the list contains strings then it will need quotes, integers in the list wouldn't have any quotes": {
                "item1": "Python 2.7.16 (default, Oct  7 2019, 17:36:04) \n[GCC 8.3.0] on linux2\nType \\"help\\", \\"copyright\\", \\"credits\\" or \\"license\\" for more information.\n>>> x = [\\"glass\\", \\"root beer\\", \\"vanilla ice cream\\", \\"straw\\"]\n>>> type(x)\n<type 'list'>\n>>> for index, ingredient in enumerate(x):\n...     print(index, ingredient)\n... \n(0, 'glass')\n(1, 'root beer')\n(2, 'vanilla ice cream')\n(3, 'straw')",
                "image": {
                    "slide9.jpg": "bottom_right"
                }
            },
            "slide_layout": 2
        },
        "Dictionaries": {
            "Dictionaries are denoted by curly braces { } and contain \\"key\\": \\"value\\" pairs. Another name for a dictionary is a Map.\nDictionaries can contain all of the previously mentioned data types including integers, strings, boolean, lists and more.": {
                "item1": ">>> x = { \\"christmas\\": \\"tree\\",\n...       \\"thanksgiving\\": \\"turkey\\",\n...       \\"halloween\\": \\"jack-o-lantern\\",\n...       \\"easter\\": \\"bunny\\"\n...     }\n>>> type(x)\n<type 'dict'>\n>>> for key, value in x.items():\n...     print(key, value)\n... \n('easter', 'bunny')\n('halloween', 'jack-o-lantern')\n('christmas', 'tree')\n('thanksgiving', 'turkey')",
                "image": {
                    "slide10.jpg": "bottom_right"
                }
            },
            "slide_layout": 2
        },
        "Hands-on Lab!": {
            "": {
                "image": {
                    "slide11.jpg": "center"
                }
            }
        }
    }
    """

    picture_position = {
        "top": {
            "left": 0.15,
            "top": 0.01,
            "width": 0.7
        },
        "bottom": {
            "left": 0.32,
            "top": 0.57,
            "width": 0.35
        },
        "right": {
            "left": 0.55,
            "top": 0.3,
            "width": 0.4
        },
        "top_right": {
            "left": 0.3,
            "top": 0.01,
            "width": 0.7
        },
        "bottom_right": {
            "left": 0.6,
            "top": 0.7,
            "width": 0.3
        },
        "top_left": {
            "left": 0.01,
            "top": 0.01,
            "width": 0.7
        },
        "bottom_left": {
            "left": 0.01,
            "top": 0.4,
            "width": 0.7
        },
        "center": {
            "left": 0.05,
            "top": 0.2,
            "width": 0.9
        }
    }

    titles = json.loads(data, object_pairs_hook=collections.OrderedDict, strict=False)
    prs = Presentation(input)

    # Create slides
    for index, (title_key, title_value) in enumerate(titles.iteritems()):

        # Set default slide layout
        default_layout = 1

        # Set slide_layout if specified
        for subtitle in title_value.keys():
            if 'slide_layout' in subtitle:
                slide_layout = title_value.get('slide_layout')
                title_value.pop('slide_layout', None)
            else:
                slide_layout = default_layout

        slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
        title = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        title.text = '{}'.format(title_key)

        # Add subtitle to slide
        for subtitle in title_value.keys():
            p = tf.add_paragraph()
            p.text = '{}'.format(subtitle)

        # Add content to slide
        for slide_index, slide_value in enumerate(title_value.values()):

            # Add paragraphs and images to slide
            for content_index, (content_key, content_value) in enumerate(slide_value.items()):

                # Add bullets
                if "item" in content_key:
                    p = tf.add_paragraph()
                    p.level = 1
                    p.font.size = Pt(16)
                    p.text = '{}'.format(content_value)

                # Add images
                if "image" in content_key:
                    for image, position in content_value.items():
                        img = imageio.imread(image)
                        picture_left = int(prs.slide_width * picture_position.get(position).get('left'))
                        picture_top = int(prs.slide_height * picture_position.get(position).get('top'))
                        picture_width = int(prs.slide_width * picture_position.get(position).get('width'))

                        picture_height = int(picture_width * img.shape[0] / img.shape[1])
                        picture = slide.shapes.add_picture(image, picture_left, picture_top, picture_width, picture_height)

    prs.save(output)


if __name__ == "__main__":
    args = parse_args()
    analyze_ppt(args.infile.name, args.outfile.name)
